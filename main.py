from http.client import HTTPException
import json
from typing import Annotated, Any, Dict
from fastapi import FastAPI, File, Query, Response, UploadFile, status
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.exceptions import HTTPException
import os
from pptx import Presentation
from python_pptx_text_replacer import TextReplacer
from pptx.chart.data import ChartData, CategoryChartData



app = FastAPI()


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(BASE_DIR, "result")
UPLOAD_DIR2 = os.path.join(BASE_DIR, "template")

def replace_chart_with_data(slide, chart_index, chart_data):
        chart_count = 0
        for shape in slide.shapes:
            if shape.has_chart:
                chart_count += 1
                if chart_count == chart_index + 1:  
                    chart = shape.chart
                    chart.replace_data(chart_data)
                    print(f"data:{chart_data}")


                    print(f"Chart with index {chart_index} found and replaced successfully on the specified slide.")
                    return
            print(f"Chart with index {chart_index} not found on the specified slide.")



@app.post('/generate')
def generate(request: Dict[Any, Any]):
    print(f"{request}")
    file = "template.pptx"

    results = request
    result = results['result']
    online_media = result['top_10_online_media']
    print_media = result['top_10_printed_media']
    perday = result['per_day_detail']
    sentiment = result['sentiment']


    print(os.path.join(UPLOAD_DIR2, file))
    prs = Presentation(os.path.join(UPLOAD_DIR2, file))

    # online media
    slide_index = 2
    chart_index_to_replace = 0
    x = list(online_media.keys())
    y = list(online_media.values())

    sorted_data = sorted(zip(x, y), key=lambda item: item[1], reverse=False)
    values = [item[1] for item in sorted_data]
    chart_data = ChartData()
    chart_data.categories = x
    chart_data.add_series('',values )
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
    
    # print media
    slide_index = 2
    chart_index_to_replace = 1
    new_chart_data = ChartData()
    x = list(print_media.keys())
    y = list(print_media.values())
    sorted_data = sorted(zip(x, y), key=lambda item: item[1], reverse=False)
    values = [item[1] for item in sorted_data]
    new_chart_data.categories = x
    new_chart_data.add_series('',values )
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_chart_data)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_chart_data)
    
    # influ online
    slide_index = 2
    chart_index_to_replace = 1
    new_chart = ChartData()
    x = list(result['top_10_online_influencer'].keys())
    y = list(result['top_10_online_influencer'].values())
    sorted_data = sorted(zip(x, y), key=lambda item: item[1], reverse=False)
    values = [item[1] for item in sorted_data]
    new_chart.categories = x
    new_chart.add_series('',values )
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_chart)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_chart_data)            
    
    # influ off
    slide_index = 4
    chart_index_to_replace = 1
    new_data = ChartData()
    x = list(result['top_10_printed_influencer'].keys())
    y = list(result['top_10_printed_influencer'].values())
    sorted_data = sorted(zip(x, y), key=lambda item: item[1], reverse=False)
    values = [item[1] for item in sorted_data]
    new_data.categories = x
    new_data.add_series('',values )
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_data)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_data)    
                
    # influ on
    slide_index = 4
    chart_index_to_replace = 0
    baru_new_chart_data = CategoryChartData()
    x = list(result['top_10_online_influencer'].keys())
    y = list(result['top_10_online_influencer'].values())
    sorted_data = sorted(zip(x, y), key=lambda item: item[1], reverse=False)
    values = [item[1] for item in sorted_data]
    baru_new_chart_data.categories = x
    baru_new_chart_data.add_series('',values )
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, baru_new_chart_data)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, baru_new_chart_data)           
            
    slide_index = 3
    chart_index_to_replace = 1
    tambah_new_chart_data = ChartData()
    x = 'positive','neutral','negative'
    y =  sentiment['positive']['percentage'],sentiment['neutral']['percentage'],sentiment['negative']['percentage']
    tambah_new_chart_data.categories = x
    tambah_new_chart_data.add_series('',y )
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, tambah_new_chart_data)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, tambah_new_chart_data)

    slide_index = 3
    chart_index_to_replace = 0
    new_new_chart_data = ChartData()
    data = result['all_days_detail']
    x = [item['text'] for item in data]
    y = [item['percentage'] for item in data]
    new_new_chart_data.categories = x
    new_new_chart_data.add_series('',y )

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_new_chart_data)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_new_chart_data)

    slide_index = 2
    chart_index_to_replace = 2

    tanggal = []
    onln = []
    pr  =  []


    for day in perday:
        onln.append(perday[day]['online'])
        pr.append(perday[day]['printed'])
        tanggal.append(day)

    data_chart_data = ChartData()
    data_chart_data.categories = tanggal
    data_chart_data.add_series('', onln)
    data_chart_data.add_series('', pr)

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, data_chart_data)
            break  

    new_filename = "result.pptx"
    SAVE_F = os.path.join(UPLOAD_DIR, new_filename)
    prs.save(SAVE_F)

    positive_topics = sentiment['positive']['topic_examples']
    negative_topics = sentiment['negative']['topic_examples']
    # neutral_topics = sentiment['neutral']['topic_examples']
    # if len(negative_topicss) > 6:
    #     negative_topics = negative_topicss
    # else:
    #     negative_topics = ''
    if tanggal:
        topics = result['per_day_detail'][tanggal[0]]
        siz = [topic['text'] for topic in topics.get('topics', [])]
    else:
        siz = []
    if len(tanggal) > 6:
        topicss = result['per_day_detail'][tanggal[6]]
        lab = [topic['text'] for topic in topicss.get('topics', [])]
    else:
        lab = []
    replace = TextReplacer(SAVE_F, slides='', tables=False, charts=False, textframes=True)
    replace.replace_text([
    ('1 – 7 April 2022', result['earliest_date']+" Sampai "+ result['latest_date']),
    ('2.251', str(result['total_online_news'])),
    ('761', str(result['total_online_media'])),
    ('135', str(result['total_printed_news'])),
    ('60', str(result['total_printed_media'])),
    ('Perluasan implementasi QRISS', siz[0]),
    ('Gangguan layanan mobile banking BCA', siz[1]),
    ('BPJPH kembangkan Sistem Informasi Halal (Sihalal) yang terintegrasi dengan penyedia uang elektronik', siz[2]),
    ('Promo belanja menggunakan kartu debit dan kredit serta dompet digitals', siz[3]),
    ('Perluasan implementasi QRIS ', lab[3]),
    ('BI dorong penggunaan transaksi nontunai selama Ramadan dan Idulfitri', lab[1]),
    ('Promo belanja menggunakan kartu debit dan kredit serta dompet digital', lab[8]),
    ('Pemerintah berencana memajaki fintech dan dompet digital', lab[9]),
    ('BI mendorong perluasan transaksi nontunai di masyarakat. ', positive_topics[3]),
    ('Kontribusi perbankan, penyedia dompet digital, dan pemerintah mendorong transaksi nontunai.', positive_topics[4]),
    ('Perbankan pastikan keamanan jaringan untuk transaksi di ATM selama Ramadan dan Idulfitri. ', positive_topics[5]),
    ('Pemerintah berencana memudahkan transaksi Pemda melalui Kartu Kredit Pemerintah Daerah (KKPD).', positive_topics[6]),
    ('Keluhan masyarakat perihal gangguan pada layanan mobile banking BCA. [link]', negative_topics[0]if isinstance(negative_topics, list) and negative_topics else ''),
    ('Terungkapnya modus skimming melalui modus pengganjal ATM di Cilacap. [link]', negative_topics[1]if isinstance(negative_topics, list) and negative_topics else ''),
    ('Terungkapnya dugaan kasus skimming nasabah BNI di Samarinda. [link] ', negative_topics[2]if isinstance(negative_topics, list) and len(negative_topics) > 2 else ''),
    ('Pencatutan identitas sebabkan kerugian berupa kesulitan pengajuan kartu kredit. [link]', negative_topics[3]if isinstance(negative_topics, list) and len(negative_topics) > 3 else ''),
    ('Keluhan soal saldo yang tidak kunjung masuk meskipun proses scan QRIS sudah berhasil. [link] ', negative_topics[4]if isinstance(negative_topics, list) and len(negative_topics) > 4 else ''),
    ('Ketimpangan penyaluran pinjaman online antara Pulau Jawa dan wilayah lainnya. [link]', negative_topics[5]if isinstance(negative_topics, list) and len(negative_topics) > 5 else ''),
])

    replace.write_presentation_to_file(SAVE_F) 

    file_path = os.path.join(UPLOAD_DIR, new_filename)
    return FileResponse(path=file_path, filename=new_filename)

